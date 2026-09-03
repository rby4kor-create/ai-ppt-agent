import os
from datetime import datetime
from lxml import etree

from pptx import Presentation as PptxDocument
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models.theme import Theme
from agents.image_agent import ImageAgent
from agents.powerpoint import text_metrics as tm
from utils.logger import get_logger

logger = get_logger(__name__)


# =============================================================================
# Layout constants - single source of truth for every coordinate used below.
# =============================================================================

PAGE_WIDTH = Inches(13.33)
PAGE_HEIGHT = Inches(7.5)

MARGIN = Inches(0.6)
GUTTER = Inches(0.35)
FOOTER_HEIGHT = Inches(0.38)

CONTENT_WIDTH = PAGE_WIDTH - 2 * MARGIN
CONTENT_TOP = Inches(0.55)
CONTENT_BOTTOM = PAGE_HEIGHT - Inches(0.55)

# Topic slide two-column grid (56% / 44% of content width, per brief)
LEFT_COLUMN_WIDTH = Emu(int(CONTENT_WIDTH * 0.56 - GUTTER / 2))
RIGHT_COLUMN_WIDTH = Emu(int(CONTENT_WIDTH * 0.44 - GUTTER / 2))
RIGHT_COLUMN_LEFT = Emu(int(MARGIN + LEFT_COLUMN_WIDTH + GUTTER))

# Weekly Signals slide - row geometry. Centralized here (rather than as
# magic numbers inside add_weekly_signals_slide) since row height is now
# computed dynamically from these plus the actual measured text height.
WS_NUMBER_COL_WIDTH = Inches(0.9)
WS_TEXT_LEFT_OFFSET = Inches(1.0)
WS_BADGE_LEFT_OFFSET = Inches(9.7)
WS_BADGE_WIDTH = Inches(1.0)
WS_HEADLINE_FONT_RANGE = (16, 13)
WS_DESC_FONT = 11
WS_ROW_GAP = Inches(0.22)          # breathing room between signal rows
WS_HEADLINE_DESC_GAP = Inches(0.06)  # gap between headline and description
WS_ROW_MAX_DESC_LINES = 2
WS_ROW_PADDING_BOTTOM = Inches(0.14)  # room for the divider rule under a row

# Topic slide vertical-flow spacing
TOPIC_TITLE_TOP_OFFSET = Inches(0.52)
TOPIC_WHY_GAP = Inches(0.28)        # gap between title bottom and "WHY IT MATTERS"
TOPIC_BULLET_TOP_GAP = Inches(0.38)  # gap between "WHY IT MATTERS" and first bullet
TOPIC_BULLET_GAP = Inches(0.14)      # gap between wrapped bullets
TOPIC_EI_RULE_GAP = Inches(0.01)     # gap between last bullet and the EI divider rule
TOPIC_EI_LABEL_GAP = Inches(0.14)    # gap between divider rule and "ENTERPRISE INTELLIGENCE"
TOPIC_STRATEGIC_LABEL_GAP = Inches(0.32)
TOPIC_STRATEGIC_TEXT_GAP = Inches(0.55)
TOPIC_RECOMMENDATION_GAP = Inches(0.16)


# =============================================================================
# Text-safety helpers (module-level - no slide/shape state needed)
# =============================================================================

def truncate_words(text, max_words):
    """Word-boundary truncation that preserves meaning better than a
    blind character cut - drops trailing filler words rather than
    cutting mid-word."""
    if not text:
        return ""
    words = str(text).split()
    if len(words) <= max_words:
        return str(text)
    return " ".join(words[:max_words]).rstrip(",.;:") + "..."


def smart_shorten_title(title, max_words=12):
    """
    Intelligently shortens a long headline instead of shrinking its font.
    Strips common preamble patterns ("Built in Fort Worth: ...") and
    otherwise truncates at a word boundary. The full original title is
    preserved separately for the references slide.
    """
    if not title:
        return ""
    if ":" in title:
        _, _, remainder = title.partition(":")
        remainder = remainder.strip()
        if len(remainder.split()) >= 5:
            title = remainder
    return truncate_words(title, max_words)


def fit_title(text, box_width_in, max_lines=3, font_range=(22, 30)):
    """
    Returns (display_text, font_pt, line_count) for a title box. Tries
    the largest font in font_range first; steps down if needed; if it
    still doesn't fit at the smallest allowed size, truncates the text
    at a word boundary rather than shrinking the font further.

    Backed by real text measurement (agents.powerpoint.text_metrics),
    not a character-count guess - the returned line_count is exact for
    the returned display_text/font_pt, so callers can position the
    next element below it without drift.
    """
    return tm.fit_text_to_box(text, box_width_in, max_lines, font_range, bold=True, step=2)


# =============================================================================


class PowerPointAgent:

    def __init__(self):
        self.output_dir = "output"
        self.output_file = os.path.join(self.output_dir, "Weekly_GenAI_Report.pptx")
        self.image_agent = ImageAgent()
        self._page_number = 0
        self._section_counter = 0

    # ------------------------------------------------------------------
    # Entry point
    # ------------------------------------------------------------------

    def generate(self, presentation_model, theme=None):
        """
        `theme` selects the palette ("Bosch Corporate" or "Modern
        Executive") for this render. Applied at the very top so every
        slide-building call below - all of which read `Theme.X` as a
        bare class attribute - picks up the chosen palette. Defaults to
        whatever palette is already active if not given, so existing
        callers that don't pass a theme keep working unchanged.
        """
        if theme:
            Theme.apply(theme)

        # Read once and stashed on the instance so every add_*_slide method
        # below can stamp the small "CW28, 2026" corner tag without every
        # call site having to thread the full presentation object through.
        self._week_label = getattr(presentation_model, "week_label", "") or ""

        os.makedirs(self.output_dir, exist_ok=True)

        prs = PptxDocument()
        prs.slide_width = PAGE_WIDTH
        prs.slide_height = PAGE_HEIGHT
        blank_layout = prs.slide_layouts[6]

        self._page_number = 0
        self._section_counter = 0
        self._recolor_theme_hyperlinks(prs)

        self.add_cover_slide(prs, blank_layout, presentation_model)
        self.add_executive_overview_slide(prs, blank_layout, presentation_model)

        if presentation_model.weekly_signals:
            self.add_weekly_signals_slide(prs, blank_layout, presentation_model)

        for section in presentation_model.sections:
            if section.get("needs_divider"):
                self.add_section_divider_slide(prs, blank_layout, section["category"], len(section["slides"]))
            for slide_model in section["slides"]:
                self.add_topic_slide(prs, blank_layout, slide_model)

        if presentation_model.strategic_takeaways:
            self.add_strategic_takeaways_slide(prs, blank_layout, presentation_model)

        # No trailing "Sources & Further Reading" slide - matching the
        # TopGenAI-CWxx-2026.pptx reference decks, every topic slide's own
        # "Source: ..." line is already a live hyperlink (see add_source),
        # so the source is available right where the claim is made instead
        # of being collected into a separate appendix at the end.

        self._validate_deck(prs)

        prs.save(self.output_file)

        logger.info("PowerPoint generated successfully.")
        logger.info(f"Saved to: {self.output_file}")

        return self.output_file

    def _recolor_theme_hyperlinks(self, prs):
        """
        PowerPoint/LibreOffice render a hyperlinked run using the slide
        master theme's <a:hlink>/<a:folHlink> scheme colors, NOT the
        run's own explicit font color - so every "Source: ..." link and
        every reference-title link was rendering in the default Office
        theme's blue/purple regardless of the run.font.color.rgb we set,
        which looked broken and off-brand against a dark premium theme.
        Patching the theme XML's hlink/folHlink to the active accent
        color fixes every hyperlink in the deck in one place.
        """
        try:
            master = prs.slide_masters[0]
            theme_part = next(
                rel.target_part for rel in master.part.rels.values()
                if "theme" in rel.reltype
            )
            root = etree.fromstring(theme_part.blob)
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            hex_color = "%02X%02X%02X" % (Theme.PRIMARY[0], Theme.PRIMARY[1], Theme.PRIMARY[2])
            for tag in ("hlink", "folHlink"):
                el = root.find(f".//a:clrScheme/a:{tag}/a:srgbClr", ns)
                if el is not None:
                    el.set("val", hex_color)
            # theme1.xml loads as a generic (non-XmlPart) Part in this
            # python-pptx version, so it exposes a plain settable `.blob`
            # rather than an `._element` tree that auto-reserializes.
            theme_part.blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
        except Exception as e:
            logger.warning(f"Could not recolor theme hyperlinks (non-fatal, links still work): {e}")

    # ------------------------------------------------------------------
    # Validation layer - runs once, after every slide has been built,
    # right before the deck is saved. Two checks:
    #   1. Boundary: does any shape leave the safe slide area?
    #   2. Collision: do any two *text-bearing* shapes' boxes overlap?
    # Both are geometric checks against the coordinates this file itself
    # set, so they're cheap and don't need to re-render anything. This
    # doesn't auto-fix problems (the dynamic-flow layout above is what
    # prevents them) - it's the safety net that surfaces a regression
    # loudly in the logs instead of shipping a silently-broken deck.
    # ------------------------------------------------------------------

    def _validate_deck(self, prs):
        total_issues = 0
        for slide_index, slide in enumerate(prs.slides, start=1):
            total_issues += self._validate_slide_bounds(slide, slide_index)
            total_issues += self._validate_slide_collisions(slide, slide_index)
        if total_issues:
            logger.warning(f"Layout validation found {total_issues} issue(s) - see warnings above.")
        else:
            logger.info("Layout validation passed: 0 boundary violations, 0 text collisions.")
        return total_issues

    def _validate_slide_bounds(self, slide, slide_index):
        issues = 0
        eps = Emu(Pt(1))  # 1pt tolerance for rounding
        for shape in slide.shapes:
            try:
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
            except Exception:
                continue
            if left is None or top is None or width is None or height is None:
                continue
            if left < -eps or top < -eps or (left + width) > PAGE_WIDTH + eps or (top + height) > PAGE_HEIGHT + eps:
                logger.warning(
                    f"[Slide {slide_index}] shape '{getattr(shape, 'name', '?')}' out of bounds: "
                    f"left={left/914400:.2f}in top={top/914400:.2f}in "
                    f"right={(left+width)/914400:.2f}in bottom={(top+height)/914400:.2f}in "
                    f"(page is {PAGE_WIDTH/914400:.2f}x{PAGE_HEIGHT/914400:.2f}in)"
                )
                issues += 1
        return issues

    def _validate_slide_collisions(self, slide, slide_index):
        boxes = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            text = tf.text.strip()
            if not text or shape.left is None:
                continue

            width, height = shape.width, shape.height
            # A short, single-line label (e.g. "GENERATED", a KPI value)
            # is very often placed in a box wider than the text itself,
            # by design, so two such boxes can sit side by side without
            # any visible collision even though their boxes technically
            # overlap. Use the actually-rendered text box (measured
            # width, capped to the shape's box) instead of the raw shape
            # box for single-line, non-wrapping content so the check
            # reflects what a viewer would actually see.
            if "\n" not in text and len(tf.paragraphs) == 1:
                try:
                    p = tf.paragraphs[0]
                    font_pt = (p.font.size.pt if p.font.size else 14)
                    bold = bool(p.font.bold)
                    measured_in = tm.measure_width_in(text, font_pt, bold=bold)
                    measured_w = Emu(int(measured_in * 914400))
                    if not tf.word_wrap or measured_w < width:
                        width = min(width, measured_w) if measured_w > 0 else width
                    line_h = Emu(int(tm.line_height_in(font_pt, p.line_spacing or 1.0) * 914400))
                    if line_h < height:
                        height = line_h
                except Exception:
                    pass

            boxes.append((text[:40], shape.left, shape.top, width, height))

        issues = 0
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                t1, l1, top1, w1, h1 = boxes[i]
                t2, l2, top2, w2, h2 = boxes[j]
                overlap_x = min(l1 + w1, l2 + w2) - max(l1, l2)
                overlap_y = min(top1 + h1, top2 + h2) - max(top1, top2)
                if overlap_x > Pt(2) and overlap_y > Pt(2):
                    # Small overlaps are normal (e.g. a bullet dot sitting
                    # just left of its text) - only flag overlaps deep
                    # enough to plausibly be visible text-on-text collision.
                    overlap_area = (overlap_x / 914400) * (overlap_y / 914400)
                    if overlap_area < 0.015:
                        continue
                    logger.warning(
                        f"[Slide {slide_index}] possible text collision: "
                        f"'{t1}' overlaps '{t2}' by {overlap_x/914400:.2f}x{overlap_y/914400:.2f}in"
                    )
                    issues += 1
        return issues

    # ------------------------------------------------------------------
    # Layout helpers
    # ------------------------------------------------------------------

    def add_rule(self, slide, left, top, width, height=Pt(2.2), color=None):
        color = color or Theme.PRIMARY
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rule.fill.solid()
        rule.fill.fore_color.rgb = color
        rule.line.fill.background()
        rule.shadow.inherit = False
        return rule

    def add_text(self, slide, left, top, width, height, text, size=Theme.BODY,
                 color=None, bold=False, align=PP_ALIGN.LEFT, font=None,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.0, word_wrap=True):
        color = color or Theme.TEXT
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = word_wrap
        tf.vertical_anchor = anchor
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.line_spacing = line_spacing
        p.font.name = font or Theme.FONT
        p.font.size = size if isinstance(size, Pt) else Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        return box

    def add_label(self, slide, left, top, width, text, color=None):
        """Small uppercase eyebrow label (category pills, section labels)."""
        return self.add_text(
            slide, left, top, width, Inches(0.28), text.upper(),
            size=Theme.SMALL, color=color or Theme.PRIMARY, bold=True,
        )

    def add_page_number(self, slide, number):
        self.add_text(
            slide, PAGE_WIDTH - Inches(0.9), PAGE_HEIGHT - FOOTER_HEIGHT,
            Inches(0.5), FOOTER_HEIGHT, str(number),
            size=Theme.SMALL, color=Theme.MUTED, align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    def add_footer(self, slide, number, label=None):
        label = label or f"Top Gen AI | AI Executive Briefing"
        self.add_text(
            slide, MARGIN, PAGE_HEIGHT - FOOTER_HEIGHT, Inches(6), FOOTER_HEIGHT,
            label, size=Theme.SMALL, color=Theme.MUTED, anchor=MSO_ANCHOR.MIDDLE,
        )
        self.add_page_number(slide, number)

    def add_corner_tag(self, slide):
        """
        Small "CW28, 2026" tag in the top-right corner of every slide -
        the running week/date reference the report is branded around
        (matches the naming convention of the TopGenAI-CWxx-2026.pptx
        reference decks). Kept deliberately tiny and muted so it reads
        as a running header, not a competing headline.
        """
        if not getattr(self, "_week_label", ""):
            return
        self.add_text(
            slide, PAGE_WIDTH - Inches(2.4), Inches(0.28), Inches(1.8), Inches(0.26),
            self._week_label.upper(), size=Pt(9), bold=True, color=Theme.MUTED,
            align=PP_ALIGN.RIGHT,
        )

    def add_source(self, slide, left, top, width, text, url=""):
        """
        Same visual treatment as before, but now hyperlinks the source
        line to the article URL when one is available - previously this
        rendered as plain, non-clickable text even though the article
        model always carried `source_link`, so "Source: OpenAI" on every
        topic slide looked like a link but did nothing when clicked.
        """
        box = slide.shapes.add_textbox(left, top, width, Inches(0.28))
        tf = box.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = truncate_words(text, 10)
        run.font.size = Theme.SMALL
        run.font.name = Theme.FONT
        run.font.color.rgb = Theme.MUTED
        if url:
            try:
                run.hyperlink.address = url
            except Exception:
                pass
        return box

    def next_page_number(self):
        self._page_number += 1
        return self._page_number

    def page_background(self, slide):
        """Fills the page with the active theme's page background - white
        for Bosch Corporate, deep navy for Modern Executive - so switching
        themes recolors every slide without touching per-slide code."""
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = Theme.BACKGROUND

    def risk_color(self, level):
        return {"Low": Theme.RISK_LOW, "Medium": Theme.RISK_MEDIUM, "High": Theme.RISK_HIGH}.get(level, Theme.RISK_MEDIUM)

    # ------------------------------------------------------------------
    # SLIDE 1 - Cover
    # ------------------------------------------------------------------

    def add_cover_slide(self, prs, layout, presentation):
        slide = prs.slides.add_slide(layout)
        self.page_background(slide)
        # No add_corner_tag() here: the cover's hero image occupies the
        # entire top-right quadrant, and the week ("CW28, 2026") is
        # already prominent inside the title text below - a second small
        # tag would just sit on top of the photo.
        self.next_page_number()

        hero_path = self.image_agent.cover_asset_path()
        img_left = Inches(8.1)
        self.image_agent.add_cropped_image(
            slide, hero_path, img_left, Inches(0), PAGE_WIDTH - img_left, PAGE_HEIGHT
        )

        left_width = Inches(7.3)
        left_width_in = left_width / 914400

        self.add_label(slide, MARGIN, Inches(0.9), left_width, "Top Gen AI")
        self.add_rule(slide, MARGIN, Inches(1.28), Inches(0.9), Pt(3))

        title_top = Inches(1.55)
        # Long report titles used to overflow a fixed 2-line box straight
        # into the subtitle below it. Same fit-then-flow approach as the
        # topic slides: shrink the font within a sensible range first,
        # then measure the *actual* wrapped height and place the subtitle
        # below that - never a fixed offset that assumes 1-2 lines.
        display_title, title_pt, _ = fit_title(presentation.title, left_width_in, max_lines=4, font_range=(34, 46))
        title_h_in = tm.text_block_height_in(display_title, title_pt, left_width_in, bold=True, line_spacing=1.03)
        self.add_text(
            slide, MARGIN, title_top, left_width, Inches(title_h_in + 0.05),
            display_title, size=Pt(title_pt), bold=True, color=Theme.TEXT, line_spacing=1.03,
        )

        subtitle_top = title_top + Inches(title_h_in) + Inches(0.22)
        self.add_text(
            slide, MARGIN, subtitle_top, left_width, Inches(0.6),
            presentation.subtitle or "AI Executive Briefing", size=Theme.SUBTITLE, color=Theme.SECONDARY_TEXT,
        )

        meta_top = Inches(6.2)
        self.add_text(slide, MARGIN, meta_top, Inches(3), Inches(0.3), "GENERATED",
                       size=Theme.SMALL, bold=True, color=Theme.MUTED)
        self.add_text(slide, MARGIN, meta_top + Inches(0.28), Inches(3), Inches(0.35),
                       presentation.generated_date, size=Theme.BODY, color=Theme.TEXT)

        self.add_text(slide, MARGIN + Inches(2.6), meta_top, Inches(3.5), Inches(0.3),
                       "PREPARED FOR", size=Theme.SMALL, bold=True, color=Theme.MUTED)
        self.add_text(slide, MARGIN + Inches(2.6), meta_top + Inches(0.28), Inches(3.5), Inches(0.35),
                       "Enterprise Leadership", size=Theme.BODY, color=Theme.TEXT)

        return slide

    # ------------------------------------------------------------------
    # SLIDE 2 - Executive Overview
    # ------------------------------------------------------------------

    def add_executive_overview_slide(self, prs, layout, presentation):
        slide = prs.slides.add_slide(layout)
        self.page_background(slide)
        self.add_corner_tag(slide)
        page = self.next_page_number()

        self.add_label(slide, MARGIN, CONTENT_TOP, Inches(4), "Weekly Briefing")
        self.add_text(
            slide, MARGIN, CONTENT_TOP + Inches(0.3), Inches(9), Inches(0.7),
            "This Week at a Glance", size=Pt(32), bold=True, color=Theme.TEXT,
        )
        self.add_rule(slide, MARGIN, CONTENT_TOP + Inches(1.05), Inches(0.7), Pt(3))

        self.add_text(
            slide, MARGIN, CONTENT_TOP + Inches(1.3), Inches(11.3), Inches(0.9),
            presentation.executive_summary, size=Theme.BODY, color=Theme.SECONDARY_TEXT, line_spacing=1.25,
        )

        stats = list(presentation.exec_stats.items())[:4]
        kpi_top = Inches(3.15)
        kpi_h = Inches(1.65)
        kpi_gap = Inches(0.3)
        kpi_w = Emu(int((CONTENT_WIDTH - kpi_gap * (len(stats) - 1)) / max(1, len(stats))))

        for i, (label, value) in enumerate(stats):
            kx = Emu(int(MARGIN + i * (kpi_w + kpi_gap)))
            border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, kx, kpi_top, kpi_w, kpi_h)
            border.fill.solid()
            border.fill.fore_color.rgb = Theme.SOFT_BG
            border.line.color.rgb = Theme.BORDER
            border.line.width = Pt(1)
            border.shadow.inherit = False

            self.add_rule(slide, kx, kpi_top, Inches(0.35), Pt(3))
            self.add_text(slide, kx + Inches(0.05), kpi_top + Inches(0.18), kpi_w - Inches(0.1), Inches(0.75),
                           str(value), size=Theme.STAT, bold=True, color=Theme.TEXT)
            self.add_text(slide, kx + Inches(0.05), kpi_top + Inches(1.15), kpi_w - Inches(0.1), Inches(0.4),
                           label.upper(), size=Theme.SMALL, bold=True, color=Theme.MUTED)

        themes_top = kpi_top + kpi_h + Inches(0.45)
        self.add_text(slide, MARGIN, themes_top, Inches(3), Inches(0.3), "TOP THEMES",
                       size=Theme.SMALL, bold=True, color=Theme.MUTED)

        pill_top = themes_top + Inches(0.35)
        px = MARGIN
        for theme_name in presentation.exec_themes[:6]:
            pill_w = Inches(0.35 + 0.11 * len(theme_name))
            if px + pill_w > MARGIN + CONTENT_WIDTH:
                break
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, pill_top, pill_w, Inches(0.42))
            pill.adjustments[0] = 0.5
            pill.fill.solid()
            pill.fill.fore_color.rgb = Theme.SOFT_BG
            pill.line.color.rgb = Theme.BORDER
            pill.line.width = Pt(0.75)
            pill.shadow.inherit = False
            tf = pill.text_frame
            tf.margin_left = tf.margin_right = Inches(0.12)
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = theme_name
            p.alignment = PP_ALIGN.CENTER
            p.font.name = Theme.FONT
            p.font.size = Theme.SMALL
            p.font.bold = True
            p.font.color.rgb = Theme.TEXT
            px = Emu(int(px + pill_w + Inches(0.18)))

        self.add_footer(slide, page)
        return slide

    # ------------------------------------------------------------------
    # SLIDE 3 - Weekly Signals
    # ------------------------------------------------------------------

    def add_weekly_signals_slide(self, prs, layout, presentation):
        slide = prs.slides.add_slide(layout)
        self.page_background(slide)
        self.add_corner_tag(slide)
        page = self.next_page_number()

        self.add_label(slide, MARGIN, CONTENT_TOP, Inches(4), "Signal Report")
        self.add_text(slide, MARGIN, CONTENT_TOP + Inches(0.3), Inches(9), Inches(0.7),
                       "Weekly Signals", size=Pt(32), bold=True, color=Theme.TEXT)
        self.add_rule(slide, MARGIN, CONTENT_TOP + Inches(1.05), Inches(0.7), Pt(3))

        row_top = CONTENT_TOP + Inches(1.4)
        footer_top = PAGE_HEIGHT - FOOTER_HEIGHT

        signals = presentation.weekly_signals[:5]
        text_left = MARGIN + WS_TEXT_LEFT_OFFSET
        text_width_in = (WS_BADGE_LEFT_OFFSET - WS_TEXT_LEFT_OFFSET - Inches(0.25)) / 914400
        desc_line_h_in = tm.line_height_in(WS_DESC_FONT, 1.15)

        # First pass: compute each row's real content (headline forced
        # to one line via smart shortening + width-fit, description
        # wrapped/truncated to WS_ROW_MAX_DESC_LINES) and its measured
        # height, BEFORE placing anything - this is what lets us detect
        # up front whether all rows fit in the available space and
        # compress gaps rather than overflowing past the footer.
        rows = []
        for signal in signals:
            headline = smart_shorten_title(signal.get("headline", ""), max_words=14)
            headline_pt = WS_HEADLINE_FONT_RANGE[0]
            for candidate_pt in range(WS_HEADLINE_FONT_RANGE[0], WS_HEADLINE_FONT_RANGE[1] - 1, -1):
                if tm.count_wrapped_lines(headline, candidate_pt, text_width_in, bold=True) == 1:
                    headline_pt = candidate_pt
                    break
                headline_pt = candidate_pt
            if tm.count_wrapped_lines(headline, headline_pt, text_width_in, bold=True) > 1:
                # Still wraps even at the smallest allowed size (very long
                # headline) - guarantee one line by trimming at a word
                # boundary rather than letting it spill into the row below.
                headline = tm.shrink_to_single_line(headline, text_width_in, headline_pt, bold=True)
            headline_h_in = tm.line_height_in(headline_pt, 1.0)

            description, _ = tm.truncate_to_lines(
                signal.get("explanation", ""), WS_DESC_FONT, text_width_in,
                max_lines=WS_ROW_MAX_DESC_LINES, bold=False,
            )
            desc_lines = tm.count_wrapped_lines(description, WS_DESC_FONT, text_width_in, bold=False)
            desc_h_in = desc_line_h_in * desc_lines

            row_content_h_in = headline_h_in + (WS_HEADLINE_DESC_GAP / 914400) + desc_h_in
            rows.append(dict(
                headline=headline, headline_pt=headline_pt, headline_h_in=headline_h_in,
                description=description, desc_h_in=desc_h_in,
                impact=signal.get("impact", "Medium"),
                content_h_in=row_content_h_in,
            ))

        # Compress gaps (never the text itself) if the natural layout
        # would run past the footer - e.g. every signal this week has a
        # 2-line description. Content is already at its guaranteed-fit
        # size from the pass above, so we only ever adjust whitespace here.
        available_in = (footer_top - Inches(0.12) - row_top) / 914400
        content_total_in = sum(r["content_h_in"] for r in rows)
        gap_count = max(1, len(rows))  # one gap+padding slot per row (incl. last, for symmetry)
        row_gap_in = WS_ROW_GAP / 914400
        row_pad_in = WS_ROW_PADDING_BOTTOM / 914400
        natural_total_in = content_total_in + gap_count * (row_gap_in + row_pad_in)
        if natural_total_in > available_in and gap_count > 0:
            slack_in = max(0.0, available_in - content_total_in)
            row_gap_in = max(0.06, slack_in / gap_count * 0.5)
            row_pad_in = max(0.05, slack_in / gap_count * 0.5)

        # Second pass: place each row using a running cursor - the next
        # row always starts at the previous row's actual measured bottom,
        # so row height is never a fixed guess.
        cursor_in = row_top / 914400
        for i, row in enumerate(rows):
            y = Emu(int(cursor_in * 914400))

            self.add_text(slide, MARGIN, y, WS_NUMBER_COL_WIDTH, Inches(0.7),
                           f"{i + 1:02d}", size=Pt(30), bold=True, color=Theme.PRIMARY)

            self.add_text(slide, text_left, y, Inches(text_width_in), Inches(row["headline_h_in"] + 0.04),
                           row["headline"], size=Pt(row["headline_pt"]), bold=True, color=Theme.TEXT)

            desc_top = y + Inches(row["headline_h_in"]) + WS_HEADLINE_DESC_GAP
            self.add_text(slide, text_left, desc_top, Inches(text_width_in), Inches(row["desc_h_in"] + 0.04),
                           row["description"], size=Pt(WS_DESC_FONT), color=Theme.SECONDARY_TEXT, line_spacing=1.15)

            impact = row["impact"]
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            MARGIN + WS_BADGE_LEFT_OFFSET, y + Inches(0.03),
                                            WS_BADGE_WIDTH, Inches(0.34))
            badge.adjustments[0] = 0.5
            badge.fill.solid()
            badge.fill.fore_color.rgb = self.risk_color(impact) if impact in ("Low", "Medium", "High") else Theme.SOFT_BG
            badge.line.fill.background()
            badge.shadow.inherit = False
            tf = badge.text_frame
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = impact
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Theme.SMALL
            p.font.bold = True
            p.font.color.rgb = Theme.WHITE
            p.font.name = Theme.FONT

            row_bottom_in = cursor_in + row["content_h_in"]
            cursor_in = row_bottom_in + row_gap_in + row_pad_in

            if i < len(rows) - 1:
                rule_y = Emu(int((row_bottom_in + row_gap_in * 0.5) * 914400))
                self.add_rule(slide, MARGIN, rule_y, CONTENT_WIDTH, Pt(0.75), color=Theme.BORDER)

        self.add_footer(slide, page)
        return slide

    # ------------------------------------------------------------------
    # Section divider
    # ------------------------------------------------------------------

    def add_section_divider_slide(self, prs, layout, category, topic_count):
        slide = prs.slides.add_slide(layout)
        self.page_background(slide)
        # No add_corner_tag() here either - same full-height hero photo on
        # the right as the cover slide.
        page = self.next_page_number()

        img_path = self.image_agent.resolve_image_path(category)
        img_left = Inches(8.3)
        self.image_agent.add_cropped_image(
            slide, img_path, img_left, Inches(0), PAGE_WIDTH - img_left, PAGE_HEIGHT
        )

        self._section_counter += 1
        section_number = self._section_counter

        self.add_rule(slide, MARGIN, Inches(2.05), Inches(0.7), Pt(3))

        self.add_text(slide, MARGIN, Inches(2.2), Inches(3), Inches(1.3),
                       f"{section_number:02d}", size=Pt(64), bold=True, color=Theme.PRIMARY)

        self.add_text(slide, MARGIN, Inches(3.5), Inches(7.5), Inches(1.3),
                       category.upper(), size=Pt(38), bold=True, color=Theme.TEXT, line_spacing=1.0)

        subtitle_map = {
            "AI Hardware": "Compute, accelerators and deployment economics",
            "AI Infrastructure": "Datacenter capacity, networking and scaling",
            "Cloud AI": "Managed platforms and infrastructure services",
            "AI Agents": "Autonomous and multi-step workflow tooling",
            "Robotics": "Physical-world manipulation and autonomy",
            "Cybersecurity": "AI safety, alignment and security posture",
            "Healthcare AI": "Clinical and healthcare-operations applications",
            "Computer Vision": "Visual recognition and inspection systems",
            "Speech AI": "Voice, transcription and audio understanding",
            "Developer AI": "AI-assisted software development tooling",
            "Generative AI": "Generative content and multimodal models",
            "Large Language Models": "Frontier and open-weight model capability",
            "General AI": "Broader signals across the AI market",
        }
        self.add_text(slide, MARGIN, Inches(4.65), Inches(6.5), Inches(0.5),
                       subtitle_map.get(category, ""), size=Theme.SUBTITLE, color=Theme.SECONDARY_TEXT)

        self.add_footer(slide, page)
        return slide

    # ------------------------------------------------------------------
    # Topic slide (the most important one)
    # ------------------------------------------------------------------

    def add_topic_slide(self, prs, layout, topic):
        slide = prs.slides.add_slide(layout)
        self.page_background(slide)
        self.add_corner_tag(slide)
        page = self.next_page_number()

        left_x = MARGIN
        left_w_in = LEFT_COLUMN_WIDTH / 914400  # EMU -> inches for the heuristic helpers

        pill_w = Inches(0.35 + 0.11 * len(topic.category))
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, CONTENT_TOP, pill_w, Inches(0.36))
        pill.adjustments[0] = 0.5
        pill.fill.solid()
        # Solid PRIMARY fill + white text (rather than a light tint with
        # dark text) keeps contrast correct on both the light and dark
        # palettes without a theme-conditional branch here.
        pill.fill.fore_color.rgb = Theme.PRIMARY
        pill.line.fill.background()
        pill.shadow.inherit = False
        tf = pill.text_frame
        tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = topic.category.upper()
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Theme.SMALL
        p.font.bold = True
        p.font.color.rgb = Theme.WHITE
        p.font.name = Theme.FONT

        title_top = CONTENT_TOP + TOPIC_TITLE_TOP_OFFSET
        display_title, title_pt, title_lines = fit_title(topic.title, left_w_in, max_lines=3, font_range=(22, 28))
        title_h_in = tm.text_block_height_in(display_title, title_pt, left_w_in, bold=True, line_spacing=1.05)
        self.add_text(slide, left_x, title_top, LEFT_COLUMN_WIDTH, Inches(title_h_in + 0.05),
                       display_title, size=Pt(title_pt), bold=True, color=Theme.TEXT, line_spacing=1.05)

        # Dynamic vertical flow starts here: every element below is
        # positioned from the *actual measured* bottom of the element
        # above it, so a 1-line and a 3-line title both flow correctly
        # with no drift accumulating down the slide.
        why_top = title_top + Inches(title_h_in) + TOPIC_WHY_GAP

        self.add_text(slide, left_x, why_top, LEFT_COLUMN_WIDTH, Inches(0.3),
                       "WHY IT MATTERS", size=Theme.SMALL, bold=True, color=Theme.PRIMARY)

        bullets = (topic.summary or [])[:3]
        bullet_top = why_top + TOPIC_BULLET_TOP_GAP
        bullet_font_pt = 13
        bullet_indent_in = 0.22
        bullet_max_width_in = left_w_in - bullet_indent_in
        footer_top = PAGE_HEIGHT - FOOTER_HEIGHT

        # The "Enterprise Intelligence" block below (divider + label +
        # Strategic + Recommendation) needs a minimum amount of room no
        # matter how long the bullets are - roughly 6 lines' worth of
        # text plus its own labels/gaps. If capping bullets at 3 lines
        # each would eat into that minimum, progressively tighten the
        # bullet cap (3 -> 2 -> 1 lines) BEFORE drawing anything, rather
        # than discovering the collision only after Recommendation has
        # already been placed on top of Source.
        ei_min_budget_in = 1.9
        chosen_cap = 3
        for cap in (3, 2, 1):
            projected_h_in = 0.0
            for bullet in bullets:
                _, wrapped = tm.truncate_to_lines(truncate_words(bullet, 40), bullet_font_pt,
                                                   bullet_max_width_in, max_lines=cap, bold=False)
                projected_h_in += tm.line_height_in(bullet_font_pt, 1.15) * len(wrapped) + 0.08 + (TOPIC_BULLET_GAP / 914400)
            projected_bottom_in = (bullet_top / 914400) + projected_h_in
            if (footer_top / 914400) - projected_bottom_in >= ei_min_budget_in:
                chosen_cap = cap
                break
            chosen_cap = cap  # keep tightening; last iteration (1) is the final fallback regardless

        cursor = bullet_top
        for bullet in bullets:
            # Guarantee-fit: truncate at a word boundary to at most
            # `chosen_cap` display lines per bullet using real wrapping,
            # instead of a blind word-count cut that could still wrap to
            # any number of lines depending on word length.
            display_bullet, wrapped = tm.truncate_to_lines(
                truncate_words(bullet, 40), bullet_font_pt, bullet_max_width_in, max_lines=chosen_cap, bold=False
            )
            row_h_in = tm.text_block_height_in(display_bullet, bullet_font_pt, bullet_max_width_in,
                                                line_spacing=1.15) + 0.08  # + small internal textbox padding
            row_h = Inches(row_h_in)

            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left_x, cursor + Inches(0.07), Inches(0.08), Inches(0.08))
            dot.fill.solid()
            dot.fill.fore_color.rgb = Theme.PRIMARY
            dot.line.fill.background()
            dot.shadow.inherit = False
            self.add_text(slide, left_x + Inches(bullet_indent_in), cursor,
                           LEFT_COLUMN_WIDTH - Inches(bullet_indent_in), row_h,
                           display_bullet, size=Pt(bullet_font_pt), color=Theme.SECONDARY_TEXT, line_spacing=1.15)

            cursor += row_h + TOPIC_BULLET_GAP

        ei_top = cursor + TOPIC_EI_RULE_GAP
        self.add_rule(slide, left_x, ei_top, LEFT_COLUMN_WIDTH, Pt(0.75), color=Theme.BORDER)
        ei_top += TOPIC_EI_LABEL_GAP
        self.add_text(slide, left_x, ei_top, LEFT_COLUMN_WIDTH, Inches(0.28),
                       "ENTERPRISE INTELLIGENCE", size=Theme.SMALL, bold=True, color=Theme.PRIMARY)

        strategic_font_pt = 11.5

        # Available vertical budget for Strategic + Recommendation before
        # the footer - computed up front so both blocks can be measured
        # against real space instead of discovering the overflow only
        # after Recommendation has already been positioned.
        budget_in = (footer_top - Inches(0.12) - (ei_top + TOPIC_STRATEGIC_LABEL_GAP)) / 914400
        # Split the budget: Strategic gets up to 2 lines, Recommendation
        # gets whatever remains (minimum 1 line) - matches the reference
        # deck's density (Strategic commonly 2 lines, Recommendation 1-2).
        strategic_line_h_in = tm.line_height_in(strategic_font_pt, 1.15)
        strategic_max_lines = max(1, min(2, int((budget_in * 0.55) // strategic_line_h_in) or 1))

        display_strategic, strategic_wrapped = tm.truncate_to_lines(
            truncate_words(topic.strategic_observation, 45), strategic_font_pt, left_w_in,
            max_lines=strategic_max_lines, bold=False,
        )
        strategic_h_in = tm.text_block_height_in(display_strategic, strategic_font_pt, left_w_in,
                                                  line_spacing=1.15) + 0.06
        strategic_h = Inches(strategic_h_in)

        self.add_text(slide, left_x, ei_top + TOPIC_STRATEGIC_LABEL_GAP, Inches(0.85), Inches(0.25),
                       "Strategic:", size=Pt(10.5), bold=True, color=Theme.MUTED)
        self.add_text(slide, left_x, ei_top + TOPIC_STRATEGIC_TEXT_GAP, LEFT_COLUMN_WIDTH, strategic_h,
                       display_strategic, size=Pt(strategic_font_pt),
                       color=Theme.TEXT, line_spacing=1.15)

        recommendation_top = ei_top + TOPIC_STRATEGIC_TEXT_GAP + strategic_h + TOPIC_RECOMMENDATION_GAP

        # Remaining real space for Recommendation, measured exactly (not
        # guessed) - guarantees Recommendation can never collide with the
        # Source line below it, however long the title/bullets/strategic
        # text above turned out to be.
        remaining_in = (footer_top - Inches(0.23) - recommendation_top) / 914400
        recommendation_max_lines = max(1, int(remaining_in // tm.line_height_in(strategic_font_pt, 1.15)))

        display_recommendation, _ = tm.truncate_to_lines(
            truncate_words(topic.recommendation, 45), strategic_font_pt, left_w_in,
            max_lines=recommendation_max_lines, bold=False,
        )
        recommendation_h_in = tm.text_block_height_in(display_recommendation, strategic_font_pt, left_w_in,
                                                       line_spacing=1.15) + 0.06
        recommendation_h = Inches(recommendation_h_in)

        self.add_text(slide, left_x, recommendation_top, Inches(1.3), Inches(0.25),
                       "Recommendation:", size=Pt(10.5), bold=True, color=Theme.MUTED)
        self.add_text(slide, left_x, recommendation_top + Inches(0.23), LEFT_COLUMN_WIDTH, recommendation_h,
                       display_recommendation, size=Pt(strategic_font_pt),
                       color=Theme.TEXT, line_spacing=1.15)

        # ---------------- RIGHT COLUMN ----------------
        img_h = Inches(2.85)
        self.image_agent.add_cropped_image(
            slide, topic.image_path, RIGHT_COLUMN_LEFT, CONTENT_TOP, RIGHT_COLUMN_WIDTH, img_h
        )

        score_top = CONTENT_TOP + img_h + Inches(0.22)
        seg_w = Emu(int(RIGHT_COLUMN_WIDTH / 3))
        seg_value_pt = 16
        seg_w_in = seg_w / 914400 - 0.08

        segs = [
            ("INNOVATION", f"{int(topic.innovation_score)}/10", Theme.TEXT),
            ("RISK", topic.risk_level, self.risk_color(topic.risk_level)),
            # READINESS is the one field of the three that's realistically
            # free text (vs. a short score/level) - guarantee-fit it to 2
            # lines so a long value can never wrap into the chip row below.
            ("READINESS", *tm.truncate_to_lines(topic.enterprise_readiness or "-", seg_value_pt, seg_w_in,
                                                 max_lines=2, bold=True)[:1], Theme.TEXT),
        ]
        # Compute the tallest of the three value strings so the chip row
        # starts right after the tallest one - never a fixed guess that a
        # 2-line READINESS value could grow past.
        max_value_lines = max(tm.count_wrapped_lines(str(v), seg_value_pt, seg_w_in, bold=True) for _, v, _ in segs)
        value_block_h_in = tm.line_height_in(seg_value_pt, 1.0) * max_value_lines

        for i, (label, value, color) in enumerate(segs):
            sx = Emu(int(RIGHT_COLUMN_LEFT + i * seg_w))
            self.add_text(slide, sx, score_top, seg_w - Inches(0.08), Inches(0.3),
                           label, size=Pt(9.5), bold=True, color=Theme.MUTED)
            self.add_text(slide, sx, score_top + Inches(0.26), seg_w - Inches(0.08), Inches(value_block_h_in + 0.06),
                           str(value), size=Pt(seg_value_pt), bold=True, color=color, line_spacing=1.0)

        chips_top = score_top + Inches(0.26) + Inches(value_block_h_in) + Inches(0.16)
        cx = RIGHT_COLUMN_LEFT
        cy = chips_top
        for tech in (topic.key_technologies or [])[:4]:
            chip_w = Inches(0.28 + 0.09 * len(tech))
            if cx + chip_w > RIGHT_COLUMN_LEFT + RIGHT_COLUMN_WIDTH:
                cx = RIGHT_COLUMN_LEFT
                cy += Inches(0.4)
            chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, chip_w, Inches(0.32))
            chip.adjustments[0] = 0.5
            chip.fill.solid()
            chip.fill.fore_color.rgb = Theme.SOFT_BG
            chip.line.color.rgb = Theme.BORDER
            chip.line.width = Pt(0.75)
            chip.shadow.inherit = False
            tf = chip.text_frame
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = tech
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(9.5)
            p.font.color.rgb = Theme.SECONDARY_TEXT
            p.font.name = Theme.FONT
            cx = Emu(int(cx + chip_w + Inches(0.12)))

        self.add_source(slide, left_x, PAGE_HEIGHT - FOOTER_HEIGHT, Inches(6),
                         f"Source: {topic.source}", url=getattr(topic, "source_link", ""))
        self.add_page_number(slide, page)

        return slide

    # ------------------------------------------------------------------
    # Strategic takeaways
    # ------------------------------------------------------------------

    def add_strategic_takeaways_slide(self, prs, layout, presentation):
        slide = prs.slides.add_slide(layout)
        self.page_background(slide)
        self.add_corner_tag(slide)
        page = self.next_page_number()

        self.add_label(slide, MARGIN, CONTENT_TOP, Inches(4), "Leadership Brief")
        self.add_text(slide, MARGIN, CONTENT_TOP + Inches(0.3), Inches(10), Inches(0.7),
                       "What This Means for Enterprise", size=Pt(30), bold=True, color=Theme.TEXT)
        self.add_rule(slide, MARGIN, CONTENT_TOP + Inches(1.0), Inches(0.7), Pt(3))

        items = presentation.strategic_takeaways[:5]
        row_top = CONTENT_TOP + Inches(1.35)
        row_h = Emu(int((CONTENT_BOTTOM - FOOTER_HEIGHT - row_top) / max(1, len(items))))

        # Build the three columns from a running x-cursor so widths and
        # gaps can never silently drift out of sync with each other
        # (a fixed-offset version of this previously overlapped the
        # implication and action columns by ~0.35in).
        col_gap = Inches(0.25)
        num_col_w = Inches(0.5)
        col_signal = Inches(2.9)
        col_impl = Inches(3.9)

        num_left = MARGIN
        signal_left = num_left + num_col_w
        impl_left = signal_left + col_signal + col_gap
        action_left = impl_left + col_impl + col_gap
        action_w = (MARGIN + CONTENT_WIDTH) - action_left

        for i, item in enumerate(items):
            y = row_top + i * row_h

            self.add_text(slide, num_left, y, num_col_w, Inches(0.4),
                           f"{i + 1:02d}", size=Pt(18), bold=True, color=Theme.PRIMARY)

            self.add_text(slide, signal_left, y, col_signal, row_h - Inches(0.15),
                           item["signal"], size=Pt(12.5), bold=True, color=Theme.TEXT, line_spacing=1.15)
            self.add_text(slide, impl_left, y, col_impl, row_h - Inches(0.15),
                           item["implication"], size=Pt(11), color=Theme.SECONDARY_TEXT, line_spacing=1.15)
            self.add_text(slide, action_left, y, action_w, row_h - Inches(0.15),
                           item["action"], size=Pt(11), color=Theme.PRIMARY_DARK, line_spacing=1.15)

            if i < len(items) - 1:
                self.add_rule(slide, MARGIN, y + row_h - Inches(0.1), CONTENT_WIDTH, Pt(0.75), color=Theme.BORDER)

        self.add_footer(slide, page)
        return slide

    # ------------------------------------------------------------------
    # References
    # ------------------------------------------------------------------

    def add_reference_slide(self, prs, layout, presentation):
        slide = prs.slides.add_slide(layout)
        self.page_background(slide)
        self.add_corner_tag(slide)
        page = self.next_page_number()

        self.add_label(slide, MARGIN, CONTENT_TOP, Inches(4), "Appendix")
        self.add_text(slide, MARGIN, CONTENT_TOP + Inches(0.3), Inches(9), Inches(0.6),
                       "Sources & Further Reading", size=Pt(28), bold=True, color=Theme.TEXT)
        self.add_rule(slide, MARGIN, CONTENT_TOP + Inches(0.95), Inches(0.7), Pt(3))

        refs = presentation.references
        col_w = Emu(int((CONTENT_WIDTH - GUTTER) / 2))
        row_h = Inches(1.0)  # generous enough for a wrapped 2-line title + link row
        rows_per_col = 5  # 5 rows x 2 cols = 10 capacity, safely fits within CONTENT_BOTTOM

        for i, ref in enumerate(refs[: rows_per_col * 2]):
            col = i // rows_per_col
            row = i % rows_per_col
            x = Emu(int(MARGIN + col * (col_w + GUTTER)))
            y = CONTENT_TOP + Inches(1.35) + row * row_h

            self.add_text(slide, x, y, col_w, Inches(0.25),
                           ref.get("source", ""), size=Pt(10.5), bold=True, color=Theme.PRIMARY)

            # Clickable row: the visible text is the clean article title
            # (not the raw URL - a bare, mid-truncated link like
            # "https://openai.com/index/gpt-5-6-la..." reads as broken
            # and gives the reader no idea what they'd be opening).
            # The full title is still word-truncated for layout safety,
            # but the hyperlink target is always the real, untruncated
            # article URL.
            link_box = slide.shapes.add_textbox(x, y + Inches(0.24), col_w, Inches(0.5))
            tf = link_box.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.1
            run = p.add_run()
            run.text = truncate_words(ref.get("title", "") or ref.get("link", ""), 12)
            run.font.size = Pt(11.5)
            run.font.color.rgb = Theme.TEXT
            run.font.name = Theme.FONT
            run.font.underline = True
            link = ref.get("link", "")
            if link:
                try:
                    run.hyperlink.address = link
                except Exception:
                    pass

        self.add_footer(slide, page)
        return slide

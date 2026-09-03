from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from models.theme import Theme


class PPTComponents:

    @staticmethod
    def draw_header(slide, title: str) -> None:

        header = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(0.45)
        )

        header.fill.solid()
        header.fill.fore_color.rgb = Theme.PRIMARY
        header.line.fill.background()

        textbox = slide.shapes.add_textbox(
            Inches(0.3),
            Inches(0.05),
            Inches(8),
            Inches(0.3)
        )

        p = textbox.text_frame.paragraphs[0]

        p.text = title

        p.font.name = Theme.FONT
        p.font.bold = True
        p.font.size = Theme.HEADER
        p.font.color.rgb = Theme.WHITE

    @staticmethod
    def draw_footer(slide, text: str) -> None:

        footer = slide.shapes.add_textbox(
            Inches(0.4),
            Inches(6.9),
            Inches(12),
            Inches(0.2)
        )

        p = footer.text_frame.paragraphs[0]

        p.text = text

        p.font.name = Theme.FONT
        p.font.size = Theme.SMALL
        p.font.color.rgb = Theme.DARK_GRAY

    @staticmethod
    def draw_footer(slide, text: str) -> None:

        from pptx.util import Pt

        title = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(0.7),
            Inches(10),
            Inches(0.5)
        )

        p = title.text_frame.paragraphs[0]

        p.text = text

        p.font.bold = True
        p.font.name = Theme.FONT
        p.font.size = Pt(26)
        p.font.color.rgb = Theme.PRIMARY

    @staticmethod
    def draw_card(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    title: str,
    value
) -> None:

        from pptx.util import Pt

        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(height)
        )

        card.fill.solid()
        card.fill.fore_color.rgb = Theme.PRIMARY_LIGHT

        card.line.color.rgb = Theme.PRIMARY

        tf = card.text_frame
        tf.clear()

        p = tf.paragraphs[0]

        p.text = str(value)

        p.alignment = PP_ALIGN.CENTER

        p.font.bold = True

        p.font.size = Pt(24)

        p.font.name = Theme.FONT

        p.font.color.rgb = Theme.PRIMARY

        p = tf.add_paragraph()

        p.text = title

        p.alignment = PP_ALIGN.CENTER

        p.font.size = Theme.BODY

        p.font.name = Theme.FONT

    @staticmethod
    def draw_section(
        slide,
        x,
        y,
        width,
        height,
        heading,
        body
    ):

        from pptx.util import Pt

        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(height)
        )

        box.fill.solid()
        box.fill.fore_color.rgb = Theme.PRIMARY_LIGHT

        box.line.color.rgb = Theme.PRIMARY

        tf = box.text_frame

        tf.clear()

        p = tf.paragraphs[0]

        p.text = heading

        p.font.bold = True

        p.font.size = Pt(16)

        p.font.name = Theme.FONT

        p.font.color.rgb = Theme.PRIMARY

        p = tf.add_paragraph()

        p.text = str(body)

        p.font.size = Theme.BODY

        p.font.name = Theme.FONT

        p.font.color.rgb = Theme.DARK_GRAY

    @staticmethod
    def draw_image_placeholder(
        slide,
        x,
        y,
        width,
        height
    ):

        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(height)
        )

        box.fill.solid()

        box.fill.fore_color.rgb = Theme.LIGHT_GRAY

        box.line.color.rgb = Theme.BORDER

        tf = box.text_frame

        tf.text = "Image"
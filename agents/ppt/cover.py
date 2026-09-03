from datetime import datetime

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from models.theme import Theme
from ppt.components import PPTComponents


class CoverSlide:

    @staticmethod
    def create(prs, report_title: str):

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Header
        PPTComponents.draw_header(
            slide,
            "AI-PPT-Agent"
        )

        # Main Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(1.2),
            Inches(11.5),
            Inches(0.8)
        )

        p = title_box.text_frame.paragraphs[0]

        p.text = report_title

        p.alignment = PP_ALIGN.CENTER

        p.font.bold = True

        p.font.size = Theme.TITLE

        p.font.name = Theme.FONT

        p.font.color.rgb = Theme.PRIMARY

        subtitle = slide.shapes.add_textbox(
            Inches(1),
            Inches(2.0),
            Inches(11),
            Inches(0.5)
        )

        p = subtitle.text_frame.paragraphs[0]

        p.text = (
            "Executive Summary of Emerging Technology Trends"
        )

        p.alignment = PP_ALIGN.CENTER

        p.font.size = Theme.SUBHEADER

        p.font.name = Theme.FONT

        p.font.color.rgb = Theme.DARK_GRAY


        generated = slide.shapes.add_textbox(
            Inches(1),
            Inches(3.2),
            Inches(11),
            Inches(0.4)
        )

        p = generated.text_frame.paragraphs[0]

        p.text = (
            f"Generated on "
            f"{datetime.now():%d %B %Y}"
        )

        p.alignment = PP_ALIGN.CENTER

        p.font.size = Theme.BODY

        p.font.name = Theme.FONT

        p.font.color.rgb = Theme.DARK_GRAY

        ai_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(4.2),
            Inches(11),
            Inches(0.4)
        )

        p = ai_box.text_frame.paragraphs[0]

        p.text = "Generated using AI-PPT-Agent"

        p.alignment = PP_ALIGN.CENTER

        p.font.bold = True

        p.font.size = Theme.BODY

        p.font.name = Theme.FONT

        p.font.color.rgb = Theme.SECONDARY

        PPTComponents.draw_image_placeholder(
            slide,
            5.4,
            4.9,
            2.5,
            1.3
        )